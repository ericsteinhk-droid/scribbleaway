"""
PowerPoint Translator GUI — translates a French .pptx slideshow to English.

Edits text in place inside the existing presentation, so images, graphics,
themes, layouts, animations and positions are untouched. Font sizes are
nudged per text box so the English text fills roughly the same space as
the French original.

Translation engine: Claude API (Anthropic). Requires an API key.
Requires: python-pptx, anthropic, Pillow
Build to .exe: pyinstaller pptx_translator.spec

Command-line mode (no GUI):
    python pptx_translator.py --cli input.pptx [--out output.pptx]
                              [--model MODEL] [--no-csv] [--mock]
"""

import csv
import hashlib
import json
import math
import os
import queue
import re
import sys
import threading
import time

# ---- optional dependencies (checked at startup) ----
try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False

try:
    import anthropic
    _ANTHROPIC_OK = True
except ImportError:
    _ANTHROPIC_OK = False

try:
    from PIL import Image, ImageTk
    _PIL_OK = True
except ImportError:
    _PIL_OK = False


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
APP_VERSION  = "v1.0"
APP_DATE     = "July 2026"
COPYRIGHT    = f"© Eric Stein, EVOQ Architecture  ·  {APP_VERSION}  ·  {APP_DATE}"

CONFIG_PATH  = os.path.join(os.path.expanduser("~"), ".pptx_translator_config.json")

MODELS = [
    ("claude-opus-4-8",  "Claude Opus 4.8 — best quality (default)"),
    ("claude-sonnet-5",  "Claude Sonnet 5 — fast, near-Opus quality"),
    ("claude-haiku-4-5", "Claude Haiku 4.5 — fastest / cheapest"),
]
DEFAULT_MODEL = "claude-opus-4-8"

BATCH_MAX_SEGMENTS = 30      # max text segments per API request
BATCH_MAX_CHARS    = 5000    # max total characters per API request

FONT_SCALE_MIN     = 0.75    # never shrink below 75 % of original size
FONT_SCALE_MAX     = 1.10    # never grow beyond 110 % of original size
FONT_SCALE_DEADBAND = 0.04   # leave size alone if within ±4 %
FONT_MIN_PT        = 6.0

DEFAULT_GLOSSARY = "EVOQ\nEVOQ Architecture\nMontréal\nQuébec"

SYSTEM_PROMPT = """You are a professional French-to-English translator working on architecture \
presentations for EVOQ Architecture, a Canadian architecture firm specializing in heritage \
conservation and contemporary design.

Rules:
- Translate French to natural, professional English suitable for client presentations.
- Preserve the exact leading and trailing whitespace of each segment.
- Keep numbers, dates, dimensions, URLs, email addresses, and file names unchanged.
- Keep proper nouns (people, firms, streets, place names, project names) unchanged unless \
they have a standard English form.
- Prefer wording of similar length to the French original when it reads naturally.
- If a segment is already in English or contains no translatable text, return it unchanged.
- Segments are fragments of slides (titles, bullets, labels, table cells) — translate each \
one independently but consistently with the others.
{glossary_block}
Return a JSON object with a "translations" array containing exactly one translated string \
per input segment, in the same order."""

TRANSLATION_SCHEMA = {
    "type": "object",
    "properties": {
        "translations": {
            "type": "array",
            "items": {"type": "string"},
        }
    },
    "required": ["translations"],
    "additionalProperties": False,
}


# ---------------------------------------------------------------------------
# Resource path (dev vs PyInstaller bundle)
# ---------------------------------------------------------------------------
def _resource_path(relative: str) -> str:
    base = getattr(sys, "_MEIPASS", os.path.dirname(os.path.abspath(__file__)))
    return os.path.join(base, relative)


# ---------------------------------------------------------------------------
# Config persistence
# ---------------------------------------------------------------------------
def load_config() -> dict:
    try:
        with open(CONFIG_PATH, "r", encoding="utf-8") as fh:
            return json.load(fh)
    except Exception:
        return {}


def save_config(cfg: dict) -> None:
    try:
        with open(CONFIG_PATH, "w", encoding="utf-8") as fh:
            json.dump(cfg, fh, indent=2)
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Translation cache (sidecar JSON so an interrupted run can resume)
# ---------------------------------------------------------------------------
class TranslationCache:
    def __init__(self, path: str, model: str, glossary: str):
        self.path = path
        self._salt = f"{model}\x1f{glossary}"
        self._data = {}
        try:
            with open(path, "r", encoding="utf-8") as fh:
                self._data = json.load(fh)
        except Exception:
            self._data = {}

    def _key(self, text: str) -> str:
        return hashlib.sha256(f"{self._salt}\x1f{text}".encode("utf-8")).hexdigest()

    def get(self, text: str):
        return self._data.get(self._key(text))

    def put(self, text: str, translation: str) -> None:
        self._data[self._key(text)] = translation

    def save(self) -> None:
        try:
            with open(self.path, "w", encoding="utf-8") as fh:
                json.dump(self._data, fh, ensure_ascii=False)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Claude translation engine
# ---------------------------------------------------------------------------
class TranslationError(Exception):
    pass


class ClaudeTranslator:
    """Batch French→English translator on the Claude Messages API."""

    def __init__(self, api_key: str, model: str = DEFAULT_MODEL, glossary: str = ""):
        if not _ANTHROPIC_OK:
            raise TranslationError("The 'anthropic' package is not installed.")
        self.model = model
        self.glossary = glossary.strip()
        # SDK retries 429 / 5xx / connection errors with exponential backoff
        self.client = anthropic.Anthropic(api_key=api_key, max_retries=4)
        terms = [t.strip() for t in self.glossary.splitlines() if t.strip()]
        if terms:
            block = ("- NEVER translate these protected terms; copy them verbatim: "
                     + "; ".join(terms) + "\n")
        else:
            block = ""
        self._system = SYSTEM_PROMPT.format(glossary_block=block)

    def translate_batch(self, texts: list, context: str = "") -> list:
        """Translate a list of segments; returns a list of the same length."""
        if not texts:
            return []
        result = self._request(texts, context)
        if result is not None:
            return result
        # Length mismatch twice in a row — fall back to one segment at a time.
        out = []
        for t in texts:
            single = self._request([t], context)
            out.append(single[0] if single else t)
        return out

    def _request(self, texts: list, context: str, _retry: bool = True):
        payload = json.dumps(
            {"context": context, "segments": texts}, ensure_ascii=False
        )
        try:
            response = self.client.messages.create(
                model=self.model,
                max_tokens=16000,
                system=self._system,
                output_config={"format": {"type": "json_schema",
                                          "schema": TRANSLATION_SCHEMA}},
                messages=[{"role": "user", "content": payload}],
            )
        except anthropic.AuthenticationError:
            raise TranslationError(
                "Invalid API key. Check your Anthropic API key and try again.")
        except anthropic.PermissionDeniedError:
            raise TranslationError(
                "This API key does not have permission to use the selected model.")
        except anthropic.NotFoundError:
            raise TranslationError(f"Model '{self.model}' was not found.")
        except anthropic.RateLimitError:
            raise TranslationError(
                "Rate limited by the API even after retries. Wait a minute and re-run "
                "— already-translated slides are cached and will not be re-sent.")
        except anthropic.APIConnectionError:
            raise TranslationError(
                "Network error talking to the Claude API. Check your internet "
                "connection and re-run — progress so far is cached.")
        except anthropic.APIStatusError as exc:
            raise TranslationError(f"Claude API error ({exc.status_code}): {exc.message}")

        if response.stop_reason == "max_tokens":
            raise TranslationError("Response truncated (max_tokens) — batch too large.")
        if response.stop_reason == "refusal":
            # Extremely unlikely for slide text; keep the French for this batch.
            return list(texts)

        text = next((b.text for b in response.content if b.type == "text"), "")
        try:
            translations = json.loads(text)["translations"]
        except Exception:
            translations = None
        if isinstance(translations, list) and len(translations) == len(texts):
            return [str(t) for t in translations]
        if _retry:
            return self._request(texts, context, _retry=False)
        return None


class MockTranslator:
    """Offline stand-in for tests (--mock). Word-level dictionary lookup."""

    _DICT = {
        "bonjour": "hello", "le": "the", "la": "the", "les": "the",
        "projet": "project", "projets": "projects", "et": "and",
        "de": "of", "du": "of the", "des": "of the", "une": "a", "un": "a",
        "architecture": "architecture", "patrimoine": "heritage",
        "bâtiment": "building", "conception": "design", "travaux": "works",
        "réhabilitation": "rehabilitation", "présentation": "presentation",
    }

    def __init__(self, *args, **kwargs):
        self.model = "mock"

    def translate_batch(self, texts: list, context: str = "") -> list:
        out = []
        for t in texts:
            words = re.split(r"(\W+)", t)
            out.append("".join(self._DICT.get(w.lower(), w) for w in words))
        return out


# ---------------------------------------------------------------------------
# PPTX text extraction / in-place editing
# ---------------------------------------------------------------------------
_LETTERS_RE = re.compile(r"[A-Za-zÀ-ÖØ-öø-ÿŒœ]")
_URL_RE     = re.compile(r"^\s*(https?://\S+|www\.\S+|\S+@\S+\.\S+)\s*$")


def needs_translation(text: str) -> bool:
    """Skip empties, pure numbers/punctuation, and bare URLs / e-mails."""
    if not text or not text.strip():
        return False
    if not _LETTERS_RE.search(text):
        return False
    if _URL_RE.match(text):
        return False
    return True


def iter_text_frames(shapes):
    """Yield every text frame under `shapes`, recursing into groups and tables."""
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from iter_text_frames(shape.shapes)
                continue
        except Exception:
            pass
        if getattr(shape, "has_text_frame", False):
            yield shape, shape.text_frame
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    yield shape, cell.text_frame


def _snap_to_space(text: str, target: int, minimum: int) -> int:
    """Nearest word boundary to `target`, never before `minimum`."""
    target = max(minimum, min(target, len(text)))
    left, right = None, None
    for i in range(target, minimum - 1, -1):
        if i == 0 or i >= len(text) or text[i] == " " or text[i - 1] == " ":
            left = i
            break
    for i in range(target, len(text) + 1):
        if i >= len(text) or text[i] == " " or (i > 0 and text[i - 1] == " "):
            right = i
            break
    if left is None:
        return right if right is not None else target
    if right is None:
        return left
    return left if (target - left) <= (right - target) else right


def split_proportional(text: str, weights: list) -> list:
    """Split `text` into len(weights) pieces sized proportionally to weights.

    Concatenation of the pieces is always exactly `text`; cuts snap to word
    boundaries so formatting changes don't land mid-word.
    """
    n = len(weights)
    if n <= 1:
        return [text]
    total = sum(weights) or 1
    pieces, prev, acc = [], 0, 0
    for w in weights[:-1]:
        acc += w
        cut = _snap_to_space(text, round(len(text) * acc / total), prev)
        pieces.append(text[prev:cut])
        prev = cut
    pieces.append(text[prev:])
    return pieces


def set_paragraph_text(paragraph, new_text: str) -> None:
    """Replace a paragraph's text, redistributing across its existing runs so
    per-run formatting (bold, italic, colour, hyperlinks) is preserved."""
    runs = paragraph.runs
    if not runs:
        return
    if len(runs) == 1:
        runs[0].text = new_text
        return
    weights = [max(len(r.text), 1) for r in runs]
    for run, piece in zip(runs, split_proportional(new_text, weights)):
        run.text = piece


def _scale_norm_autofit(text_frame, ratio: float) -> bool:
    """Scale an existing <a:normAutofit fontScale> (used when runs inherit
    their size from the layout/master and carry no explicit size)."""
    try:
        bodyPr = text_frame._txBody.find(qn("a:bodyPr"))
        if bodyPr is None:
            return False
        autofit = bodyPr.find(qn("a:normAutofit"))
        if autofit is None:
            return False
        raw = autofit.get("fontScale")
        current = 100000 if raw is None else (
            int(float(raw.rstrip("%")) * 1000) if raw.endswith("%") else int(raw))
        scaled = int(max(25000, min(100000, current * ratio)))
        autofit.set("fontScale", str(scaled))
        return True
    except Exception:
        return False


def resize_text_frame(text_frame, orig_chars: int, new_chars: int) -> float:
    """Nudge font sizes so the new text fills roughly the original space.

    For a fixed-width box, occupied height scales ~ chars × size², so the
    size correction is sqrt(orig/new), clamped to a conservative range.
    Returns the ratio applied (1.0 = untouched).
    """
    if orig_chars <= 0 or new_chars <= 0:
        return 1.0
    ratio = math.sqrt(orig_chars / new_chars)
    ratio = max(FONT_SCALE_MIN, min(FONT_SCALE_MAX, ratio))
    if abs(ratio - 1.0) <= FONT_SCALE_DEADBAND:
        return 1.0
    changed = False
    for para in text_frame.paragraphs:
        for run in para.runs:
            size = run.font.size
            if size is not None:
                new_pt = max(FONT_MIN_PT, round(size.pt * ratio * 2) / 2)
                run.font.size = Pt(new_pt)
                changed = True
    if not changed:
        # No explicit sizes anywhere — sizes are inherited. Adjust the
        # frame's autofit scale instead, if it has one; otherwise leave
        # PowerPoint's own autofit to handle it.
        if not _scale_norm_autofit(text_frame, ratio):
            return 1.0
    return ratio


# ---------------------------------------------------------------------------
# Translation worker
# ---------------------------------------------------------------------------
def default_output_path(in_path: str) -> str:
    stem, ext = os.path.splitext(in_path)
    return f"{stem}_EN{ext}"


def _chunk(items: list, max_items: int, max_chars: int, key=len):
    batch, chars = [], 0
    for item in items:
        n = key(item)
        if batch and (len(batch) >= max_items or chars + n > max_chars):
            yield batch
            batch, chars = [], 0
        batch.append(item)
        chars += n
    if batch:
        yield batch


def translate_pptx(in_path, out_path, translator, use_cache=True,
                   translate_tables=True, resize_fonts=True, write_csv=True,
                   progress_cb=None, cancel_event=None):
    """Translate `in_path` into a new file at `out_path`. Never touches the
    original. Returns a stats dict."""

    def report(msg, current=None, total=None):
        if progress_cb:
            progress_cb(msg, current, total)

    prs = Presentation(in_path)
    slides = list(prs.slides)

    cache_path = os.path.splitext(in_path)[0] + ".translation_cache.json"
    cache = TranslationCache(cache_path, translator.model,
                             getattr(translator, "glossary", ""))

    # ---- collect: one entry per paragraph that needs translating ----
    # entry = (slide_no, shape_name, text_frame, paragraph, french_text)
    entries = []
    for idx, slide in enumerate(slides, 1):
        for shape, tf in iter_text_frames(slide.shapes):
            if not translate_tables and getattr(shape, "has_table", False):
                continue
            for para in tf.paragraphs:
                text = "".join(r.text for r in para.runs)
                if needs_translation(text):
                    entries.append((idx, shape.name, tf, para, text))

    report(f"Found {len(entries)} text segments on {len(slides)} slides.",
           0, len(entries))

    # ---- translate slide by slide (whole-slide context, cached segments skipped)
    translations = {}          # french text -> english text
    done = 0
    api_calls = 0
    from_cache = 0
    by_slide = {}
    for entry in entries:
        by_slide.setdefault(entry[0], []).append(entry)

    for slide_no in sorted(by_slide):
        if cancel_event is not None and cancel_event.is_set():
            cache.save()
            raise TranslationError("Cancelled by user.")
        pending = []
        for _, _, _, _, text in by_slide[slide_no]:
            if text in translations:
                continue
            cached = cache.get(text) if use_cache else None
            if cached is not None:
                translations[text] = cached
                from_cache += 1
            elif text not in pending:
                pending.append(text)
        for batch in _chunk(pending, BATCH_MAX_SEGMENTS, BATCH_MAX_CHARS):
            if cancel_event is not None and cancel_event.is_set():
                cache.save()
                raise TranslationError("Cancelled by user.")
            context = (f"Slide {slide_no} of {len(slides)} in a French "
                       f"architecture presentation.")
            results = translator.translate_batch(batch, context)
            api_calls += 1
            for fr, en in zip(batch, results):
                translations[fr] = en
                cache.put(fr, en)
            cache.save()
        done += len(by_slide[slide_no])
        report(f"Slide {slide_no}/{len(slides)} translated.", done, len(entries))

    # ---- apply: rewrite runs in place, then resize per text frame ----
    frames = {}                # id(tf) -> (tf, orig_chars, new_chars)
    review_rows = []
    for slide_no, shape_name, tf, para, fr in entries:
        en = translations.get(fr, fr)
        set_paragraph_text(para, en)
        review_rows.append((slide_no, shape_name, fr, en))
        key = id(tf)
        _, o, n = frames.get(key, (tf, 0, 0))
        frames[key] = (tf, o + len(fr.strip()), n + len(en.strip()))

    resized = 0
    if resize_fonts:
        for tf, orig_chars, new_chars in frames.values():
            if resize_text_frame(tf, orig_chars, new_chars) != 1.0:
                resized += 1

    # ---- save atomically, then verify the output opens ----
    tmp_path = out_path + ".tmp"
    prs.save(tmp_path)
    os.replace(tmp_path, out_path)
    Presentation(out_path)     # raises if the file is corrupt

    csv_path = None
    if write_csv:
        csv_path = os.path.splitext(out_path)[0] + "_review.csv"
        with open(csv_path, "w", newline="", encoding="utf-8-sig") as fh:
            writer = csv.writer(fh)
            writer.writerow(["Slide", "Shape", "French", "English"])
            writer.writerows(review_rows)

    report("Done.", len(entries), len(entries))
    return {
        "slides": len(slides),
        "segments": len(entries),
        "unique": len(translations),
        "api_calls": api_calls,
        "from_cache": from_cache,
        "frames_resized": resized,
        "out_path": out_path,
        "csv_path": csv_path,
    }


# ---------------------------------------------------------------------------
# GUI
# ---------------------------------------------------------------------------
def run_gui():
    import tkinter as tk
    from tkinter import filedialog, messagebox, ttk

    ACCENT = "#1a5276"

    class TranslatorApp(tk.Tk):
        def __init__(self):
            super().__init__()
            self.title("EVOQ PowerPoint Translator — French → English")
            self.geometry("860x780")
            self.minsize(760, 640)
            self.configure(bg="#ffffff")
            self._queue = queue.Queue()
            self._cancel = threading.Event()
            self._worker = None
            self._cfg = load_config()
            self._build_ui()
            self.protocol("WM_DELETE_WINDOW", self._on_close)

        # ---- UI construction -------------------------------------------
        def _build_ui(self):
            # Logo header
            hbar = tk.Frame(self, bg="#ffffff")
            hbar.pack(fill="x")
            if _PIL_OK:
                try:
                    img = Image.open(_resource_path("evoq_logo.png")).convert("RGBA")
                    scale = 46 / img.height
                    img = img.resize((int(img.width * scale), 46))
                    self._logo = ImageTk.PhotoImage(img)
                    tk.Label(hbar, image=self._logo, bg="#ffffff").pack(
                        side="left", padx=14, pady=9)
                except Exception:
                    pass
            tk.Label(hbar, text="PowerPoint Translator  ·  Français → English",
                     font=("Segoe UI", 15, "bold"), fg=ACCENT,
                     bg="#ffffff").pack(side="left", padx=6)
            tk.Frame(self, height=1, bg="#d5d8dc").pack(fill="x")

            body = tk.Frame(self, bg="#ffffff")
            body.pack(fill="both", expand=True, padx=16, pady=8)
            body.columnconfigure(1, weight=1)

            # File selection
            tk.Label(body, text="PowerPoint file (.pptx):", bg="#ffffff",
                     font=("Segoe UI", 10)).grid(row=0, column=0, sticky="w", pady=4)
            self._file_var = tk.StringVar()
            tk.Entry(body, textvariable=self._file_var,
                     font=("Segoe UI", 10)).grid(row=0, column=1, sticky="ew",
                                                 padx=8, pady=4)
            tk.Button(body, text="Browse…", command=self._browse,
                      bg=ACCENT, fg="#ffffff", relief="flat", padx=12,
                      cursor="hand2").grid(row=0, column=2, pady=4)

            # API key
            tk.Label(body, text="Anthropic API key:", bg="#ffffff",
                     font=("Segoe UI", 10)).grid(row=1, column=0, sticky="w", pady=4)
            self._key_var = tk.StringVar(value=self._cfg.get("api_key", ""))
            self._key_entry = tk.Entry(body, textvariable=self._key_var, show="•",
                                       font=("Segoe UI", 10))
            self._key_entry.grid(row=1, column=1, sticky="ew", padx=8, pady=4)
            self._show_var = tk.BooleanVar(value=False)
            tk.Checkbutton(body, text="Show", variable=self._show_var,
                           bg="#ffffff", command=self._toggle_key).grid(
                row=1, column=2, pady=4)

            # Model
            tk.Label(body, text="Model:", bg="#ffffff",
                     font=("Segoe UI", 10)).grid(row=2, column=0, sticky="w", pady=4)
            labels = [label for _, label in MODELS]
            self._model_box = ttk.Combobox(body, values=labels, state="readonly",
                                           font=("Segoe UI", 10))
            saved_model = self._cfg.get("model", DEFAULT_MODEL)
            index = next((i for i, (m, _) in enumerate(MODELS) if m == saved_model), 0)
            self._model_box.current(index)
            self._model_box.grid(row=2, column=1, sticky="ew", padx=8, pady=4)

            # Glossary
            tk.Label(body, text="Do NOT translate\n(one term per line):",
                     bg="#ffffff", justify="left",
                     font=("Segoe UI", 10)).grid(row=3, column=0, sticky="nw", pady=4)
            self._glossary = tk.Text(body, height=5, font=("Segoe UI", 10),
                                     relief="solid", borderwidth=1)
            self._glossary.insert("1.0", self._cfg.get("glossary", DEFAULT_GLOSSARY))
            self._glossary.grid(row=3, column=1, columnspan=2, sticky="ew",
                                padx=8, pady=4)

            # Options
            opts = tk.Frame(body, bg="#ffffff")
            opts.grid(row=4, column=0, columnspan=3, sticky="w", pady=(6, 2))
            self._opt_tables = tk.BooleanVar(value=self._cfg.get("tables", True))
            self._opt_resize = tk.BooleanVar(value=self._cfg.get("resize", True))
            self._opt_csv    = tk.BooleanVar(value=self._cfg.get("csv", True))
            self._opt_cache  = tk.BooleanVar(value=self._cfg.get("cache", True))
            for text, var in (
                ("Translate tables", self._opt_tables),
                ("Resize fonts to fill original space", self._opt_resize),
                ("Bilingual review CSV", self._opt_csv),
                ("Use translation cache (resume after errors)", self._opt_cache),
            ):
                tk.Checkbutton(opts, text=text, variable=var,
                               bg="#ffffff", font=("Segoe UI", 9)).pack(
                    side="left", padx=(0, 14))

            # Action row
            action = tk.Frame(body, bg="#ffffff")
            action.grid(row=5, column=0, columnspan=3, sticky="ew", pady=8)
            self._go_btn = tk.Button(action, text="  Translate  ",
                                     command=self._start,
                                     bg="#1e8449", fg="#ffffff", relief="flat",
                                     font=("Segoe UI", 11, "bold"),
                                     padx=16, pady=4, cursor="hand2")
            self._go_btn.pack(side="left")
            self._cancel_btn = tk.Button(action, text="Cancel",
                                         command=self._cancel_run,
                                         state="disabled", relief="flat",
                                         bg="#c0392b", fg="#ffffff",
                                         padx=12, pady=4, cursor="hand2")
            self._cancel_btn.pack(side="left", padx=10)
            self._status_var = tk.StringVar(value="Select a .pptx file to begin.")
            tk.Label(action, textvariable=self._status_var, bg="#ffffff",
                     fg="#566573", font=("Segoe UI", 9)).pack(side="left", padx=8)

            self._progress = ttk.Progressbar(body, mode="determinate")
            self._progress.grid(row=6, column=0, columnspan=3, sticky="ew", pady=2)

            # Log
            self._log = tk.Text(body, height=14, state="disabled",
                                font=("Consolas", 9), bg="#f8f9f9",
                                relief="solid", borderwidth=1)
            self._log.grid(row=7, column=0, columnspan=3, sticky="nsew", pady=6)
            body.rowconfigure(7, weight=1)

            tk.Label(self, text=COPYRIGHT, bg="#ffffff", fg="#95a5a6",
                     font=("Segoe UI", 8)).pack(side="bottom", pady=4)

        # ---- helpers ----------------------------------------------------
        def _toggle_key(self):
            self._key_entry.config(show="" if self._show_var.get() else "•")

        def _browse(self):
            path = filedialog.askopenfilename(
                title="Select a French PowerPoint file",
                filetypes=[("PowerPoint", "*.pptx"), ("All files", "*.*")])
            if path:
                self._file_var.set(path)

        def _log_line(self, text):
            self._log.config(state="normal")
            self._log.insert("end", text + "\n")
            self._log.see("end")
            self._log.config(state="disabled")

        def _selected_model(self):
            return MODELS[self._model_box.current()][0]

        def _save_cfg(self):
            self._cfg.update({
                "api_key": self._key_var.get().strip(),
                "model": self._selected_model(),
                "glossary": self._glossary.get("1.0", "end").strip(),
                "tables": self._opt_tables.get(),
                "resize": self._opt_resize.get(),
                "csv": self._opt_csv.get(),
                "cache": self._opt_cache.get(),
            })
            save_config(self._cfg)

        def _on_close(self):
            self._save_cfg()
            self._cancel.set()
            self.destroy()

        # ---- run --------------------------------------------------------
        def _start(self):
            in_path = self._file_var.get().strip()
            key = self._key_var.get().strip()
            if not in_path or not os.path.isfile(in_path):
                messagebox.showwarning("No file", "Select a .pptx file first.")
                return
            if not in_path.lower().endswith(".pptx"):
                messagebox.showwarning(
                    "Unsupported format",
                    "Only .pptx files are supported. Open the file in PowerPoint "
                    "and save it as .pptx first.")
                return
            if not key:
                messagebox.showwarning(
                    "No API key",
                    "Enter your Anthropic API key (from console.anthropic.com).")
                return
            self._save_cfg()

            out_path = default_output_path(in_path)
            self._cancel.clear()
            self._go_btn.config(state="disabled")
            self._cancel_btn.config(state="normal")
            self._progress["value"] = 0
            self._log_line(f"Translating: {os.path.basename(in_path)}")
            self._log_line(f"Output:      {os.path.basename(out_path)}")

            translator_args = dict(
                api_key=key, model=self._selected_model(),
                glossary=self._glossary.get("1.0", "end"))
            options = dict(
                use_cache=self._opt_cache.get(),
                translate_tables=self._opt_tables.get(),
                resize_fonts=self._opt_resize.get(),
                write_csv=self._opt_csv.get())

            def progress_cb(msg, current, total):
                self._queue.put(("progress", msg, current, total))

            def worker():
                try:
                    translator = ClaudeTranslator(**translator_args)
                    stats = translate_pptx(in_path, out_path, translator,
                                           progress_cb=progress_cb,
                                           cancel_event=self._cancel, **options)
                    self._queue.put(("done", stats))
                except TranslationError as exc:
                    self._queue.put(("error", str(exc)))
                except Exception as exc:
                    self._queue.put(("error", f"Unexpected error: {exc!r}"))

            self._worker = threading.Thread(target=worker, daemon=True)
            self._worker.start()
            self.after(100, self._poll)

        def _cancel_run(self):
            self._cancel.set()
            self._status_var.set("Cancelling…")

        def _poll(self):
            try:
                while True:
                    item = self._queue.get_nowait()
                    if item[0] == "progress":
                        _, msg, current, total = item
                        self._status_var.set(msg)
                        self._log_line(msg)
                        if total:
                            self._progress["maximum"] = total
                            self._progress["value"] = current or 0
                    elif item[0] == "done":
                        self._finish(item[1])
                        return
                    elif item[0] == "error":
                        self._fail(item[1])
                        return
            except queue.Empty:
                pass
            if self._worker and self._worker.is_alive():
                self.after(100, self._poll)

        def _finish(self, stats):
            self._go_btn.config(state="normal")
            self._cancel_btn.config(state="disabled")
            self._status_var.set("Done.")
            self._log_line(
                f"Done — {stats['segments']} segments on {stats['slides']} slides "
                f"({stats['from_cache']} from cache, {stats['api_calls']} API calls, "
                f"{stats['frames_resized']} text boxes resized).")
            self._log_line(f"Saved: {stats['out_path']}")
            if stats.get("csv_path"):
                self._log_line(f"Review CSV: {stats['csv_path']}")
            messagebox.showinfo(
                "Translation complete",
                f"Translated {stats['segments']} text segments on "
                f"{stats['slides']} slides.\n\nSaved to:\n{stats['out_path']}")

        def _fail(self, message):
            self._go_btn.config(state="normal")
            self._cancel_btn.config(state="disabled")
            self._status_var.set("Failed.")
            self._log_line(f"ERROR: {message}")
            messagebox.showerror("Translation failed", message)

    missing = []
    if not _PPTX_OK:
        missing.append("python-pptx")
    if not _ANTHROPIC_OK:
        missing.append("anthropic")
    if missing:
        import tkinter as tk
        from tkinter import messagebox
        root = tk.Tk()
        root.withdraw()
        messagebox.showerror(
            "Missing dependency",
            "Required packages are not installed:\n  " + "\n  ".join(missing)
            + "\n\nInstall with:  pip install " + " ".join(missing))
        return
    TranslatorApp().mainloop()


# ---------------------------------------------------------------------------
# CLI mode (headless)
# ---------------------------------------------------------------------------
def run_cli(argv):
    import argparse
    parser = argparse.ArgumentParser(prog="pptx_translator --cli")
    parser.add_argument("input")
    parser.add_argument("--out", default=None)
    parser.add_argument("--model", default=DEFAULT_MODEL)
    parser.add_argument("--glossary", default=DEFAULT_GLOSSARY)
    parser.add_argument("--no-csv", action="store_true")
    parser.add_argument("--no-resize", action="store_true")
    parser.add_argument("--no-cache", action="store_true")
    parser.add_argument("--mock", action="store_true",
                        help="offline dictionary engine (testing only)")
    args = parser.parse_args(argv)

    out_path = args.out or default_output_path(args.input)
    if args.mock:
        translator = MockTranslator()
    else:
        key = os.environ.get("ANTHROPIC_API_KEY") or load_config().get("api_key", "")
        if not key:
            sys.exit("Set ANTHROPIC_API_KEY or run the GUI once to save a key.")
        translator = ClaudeTranslator(key, args.model, args.glossary)

    def progress(msg, current, total):
        print(msg, flush=True)

    t0 = time.time()
    stats = translate_pptx(
        args.input, out_path, translator,
        use_cache=not args.no_cache,
        resize_fonts=not args.no_resize,
        write_csv=not args.no_csv,
        progress_cb=progress)
    print(f"\n{stats['segments']} segments / {stats['slides']} slides "
          f"in {time.time() - t0:.1f}s  ->  {stats['out_path']}")
    if stats.get("csv_path"):
        print(f"Review CSV: {stats['csv_path']}")
    return stats


if __name__ == "__main__":
    if "--cli" in sys.argv:
        argv = [a for a in sys.argv[1:] if a != "--cli"]
        run_cli(argv)
    else:
        run_gui()
